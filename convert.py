import os
from wand.image import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN


class Convert():

    def __init__(self):
        print("=========PYTHON ASSIGNMENT=========")
        
        self.input_dir = input("Enter directory tto fetch image from : ")
        self.input_dir = self.input_dir.replace("\\","/")
        self.input_dir = self.input_dir+"/"
        print ('Fetching from : ' + self.input_dir)

        self.out_dir = input('Enter output directory name : ')
        output_dir = ''.join(char for char in self.out_dir if char.isalnum())
        self.out_dir = output_dir+"/"
        print(self.out_dir)

        global output_path
        output_path = os.path.join(self.input_dir,self.out_dir)

        if not os.path.exists(self.out_dir):
            os.mkdir(output_path)

    def apply_logo(self):
        print("Files will be stored in : "+output_path)
        ext = '.jpg'

        with Image(filename='nike_black.png') as img_logo:
            img_logo.resize(800,300)

            for images in os.listdir(self.input_dir):
                if images.endswith(ext):
                    print(images)

                    with Image(filename=self.input_dir + images) as img:
                        img.watermark(img_logo,0,40,40)
                        img.save(filename=output_path + "new" + images)
                        print("logo applied")
                    
                else:
                    print('No matching extention file found !!')
                    

    def make_ppt(self):
        print("Presentation will be stored in :"+output_path)
        ext= '.jpg'

        pr1 = Presentation()

        slide1_register = pr1.slide_layouts[0]
        slide1 = pr1.slides.add_slide(slide1_register)

        title = slide1.shapes.title
        subtitle = slide1.placeholders[1]

        title.text = "Python Assignment"
        subtitle.text = "Created by : Mohd Aadil Shafi"


        for images in os.listdir(self.input_dir):
            if images.endswith(ext):
                
                #Content Slide
                slide2_register = pr1.slide_layouts[6]
                slide2 = pr1.slides.add_slide(slide2_register)
                
                #First Heading
                left = Inches(3.5)
                top = Inches(0)
                width = height = Inches(1)

                txBox = slide2.shapes.add_textbox(left,top,width,height)
                tf = txBox.text_frame
                # tf.text= "Python"
                t_paragrph = tf.add_paragraph()
                t_paragrph.text = "Python pptx"
                t_paragrph.font.bold = True
                t_paragrph.font.name = 'Bell MT'
                t_paragrph.font.size = Pt(40)
                # t_paragrph = PP_ALIGN.CENTER

                #Second Headign
                top2 = Inches(1)

                txBox2 = slide2.shapes.add_textbox(left, top2,width,height)
                tf2 = txBox2.text_frame

                t_paragrph2 = tf2.add_paragraph()
                t_paragrph2.text = "Assignment"
                t_paragrph2.font.name = 'Bell MT'
                t_paragrph2.font.size = Pt(30)

                #images
                from_top = Inches(2)

                add_img = slide2.shapes.add_picture(self.input_dir+images,left,from_top,Inches(3))

        pr1.save(output_path + "Aadil-Assignment.pptx")  
        print("Presentation saved")       




first = Convert()
# first.apply_logo()
first.make_ppt()